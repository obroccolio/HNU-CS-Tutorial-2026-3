#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成PPT：矩阵乘法性能优化 - 问题d和e分析
风格：Octo Code 深色开发者设计系统
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Octo Code 色彩系统
COLORS = {
    'void_black': RGBColor(13, 17, 23),         # #0D1117
    'surface': RGBColor(22, 27, 34),            # #161B22
    'elevated': RGBColor(28, 33, 40),           # #1C2128
    'modal': RGBColor(33, 38, 45),              # #21262D
    'border': RGBColor(48, 54, 61),             # #30363D
    'muted_border': RGBColor(72, 79, 88),       # #484F58
    'text_primary': RGBColor(230, 237, 243),    # #E6EDF3
    'text_secondary': RGBColor(139, 148, 158),  # #8B949E
    'mona_blue': RGBColor(47, 129, 247),        # #2F81F7
    'mona_blue_hover': RGBColor(56, 139, 253),  # #388BFD
    'growth_green': RGBColor(35, 134, 54),      # #238636
    'success': RGBColor(63, 185, 80),           # #3FB950
    'warning': RGBColor(210, 153, 34),          # #D29922
    'error': RGBColor(248, 81, 73),             # #F85149
}

def add_title_slide(prs, title, subtitle):
    """添加标题页 - Octo Code 风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['void_black']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_primary']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Inter'
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['mona_blue']
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Inter'

def add_content_slide(prs, title, content_items):
    """添加内容页 - Octo Code 风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['void_black']
    
    # 标题栏
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['mona_blue']
    p.font.name = 'Inter'
    
    # 添加分割线
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0))
    line.line.color.rgb = COLORS['border']
    line.line.width = Pt(1)
    
    # 内容
    content_top = 1.5
    for i, item in enumerate(content_items):
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(content_top + i * 0.68), Inches(8.4), Inches(0.65))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        p = text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text_primary']
        p.font.name = 'Inter'
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.level = 0

def add_table_slide(prs, title, table_data):
    """添加表格页 - Octo Code 风格"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['void_black']
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['mona_blue']
    p.font.name = 'Inter'
    
    # 表格
    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(5.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # 填充表格
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            
            # 格式化
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.name = 'Inter'
            
            if i == 0:  # 表头
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = COLORS['text_primary']
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['surface']
            else:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = COLORS['text_primary']
                
                # 交替行背景
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS['elevated']
            
            paragraph.alignment = PP_ALIGN.CENTER

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 第1页：标题页
add_title_slide(prs, "矩阵乘法性能优化", "问题 d 和 e 分析")

# 第2页：问题概述
add_content_slide(prs, "问题概述", [
    "背景：矩阵乘法是计算密集型操作，性能与内存访问模式密切相关",
    "",
    "问题 d：分析矩阵转置对性能的影响",
    "    • 为什么转置 B 矩阵能提升性能？",
    "    • 不同循环顺序下转置的效果是否一致？",
    "",
    "问题 e：矩阵规模对缓存性能的影响",
    "    • 为什么 2 的幂次（如 1024×1024）性能骤降？",
    "    • 缓存冲突如何导致性能灾难？"
])

# 第3页：问题d - 原理
add_content_slide(prs, "问题 d：矩阵转置的妙处", [
    "基本矩阵乘法：C[i][j] = Σ A[i][k] × B[k][j]",
    "",
    "原始问题：B 矩阵的访问模式",
    "    当 k 变化时，访问 B[k][j] 的地址间隔为 N（列宽）",
    "    → 列优先访问，步长大，缓存命中率低",
    "",
    "转置的解决方案：",
    "    先转置 B → B_T，修改访问为 B_T[j][k]",
    "    当 k 变化时，访问相邻地址",
    "    → 行优先访问，步长为 1，缓存命中率高",
    "",
    "性价比：转置开销 ≈ 0.1% vs 性能收益 ≈ 30-50%"
])

# 第4页：问题d - 访问模式分析
add_table_slide(prs, "问题 d：6种循环顺序的访问模式", [
    ["循环顺序", "A访问模式", "B访问模式", "C访问模式", "性能等级"],
    ["ijk", "行优先(+1)", "列优先(+N)", "固定", "中等"],
    ["jik", "行优先(+1)", "列优先(+N)", "列优先(+1)", "较优"],
    ["kij", "固定", "行优先(+1)", "行优先(+1)", "最优"],
    ["ikj", "行优先(+K)", "列优先(+N)", "行优先(+1)", "较差"],
    ["jki", "固定", "列优先(+1)", "行优先(+N)", "较差"],
    ["kji", "行优先(+1)", "固定", "行优先(+N)", "较差"]
])

# 第5页：问题d - 性能对比
add_table_slide(prs, "问题 d：转置效果对比（-O2优化）", [
    ["循环顺序", "不转置(ms)", "转置后(ms)", "加速比", "性能提升"],
    ["ijk", "783.06", "457.40", "1.71x", "+41.59%"],
    ["jik", "767.27", "452.72", "1.70x", "+41.00%"],
    ["kij", "187.27", "387.55", "0.48x", "-51.59%"],
    ["ikj", "191.46", "321.55", "0.60x", "-40.56%"],
    ["jki", "659.78", "553.88", "1.19x", "+16.06%"],
    ["kji", "570.46", "617.21", "0.92x", "-7.87%"]
])

# 第6页：问题d - 关键发现
add_content_slide(prs, "问题 d：关键发现", [
    "转置获益的情况：",
    "    • ijk、jik 顺序：B 原本是列优先 → 转置后获得 40%+ 性能提升",
    "    • jki 顺序：获得 10-17% 性能提升",
    "",
    "转置反而减速的情况：",
    "    • kij、ikj 顺序：B 原本已是行优先 → 转置后反而变成列优先",
    "    • 转置摧毁了已有的优化结构，性能下降 40-50%",
    "",
    "结论：转置效果强烈依赖于原始循环顺序的访问模式"
])

# 第7页：问题e - 缓存冲突问题
add_content_slide(prs, "问题 e：Cache 冲突陷阱", [
    "现象：矩阵规模为 2 的幂时性能骤降",
    "    • 1023×1023：0.74x (正常)",
    "    • 1024×1024：0.24x (大幅下降！)",
    "    • 1025×1025：0.78x (恢复正常)",
    "",
    "根本原因：Cache 组冲突",
    "    CPU L2/L3 Cache 是组相联结构",
    "    组索引计算：(物理地址 >> 6) & (组数-1)",
    "",
    "当矩阵宽度是 2 的幂时：",
    "    行间地址差 = 1024×8 = 8192 = 2^13",
    "    多行映射到同一 Cache 组 → 频繁驱逐 → 缓存命中率崩溃"
])

# 第8页：问题e - 详细机制
add_content_slide(prs, "问题 e：Cache 冲突机制", [
    "Cache 结构（以M1为例）：",
    "    L2 Cache = 256KB = 512组 × 8路 × 64B行",
    "",
    "地址到组的映射（以1024×1024矩阵为例）：",
    "    行0 起始地址：0x0      → 组索引 0x000",
    "    行1 起始地址：0x2000   → 组索引 0x080",
    "    行2 起始地址：0x4000   → 组索引 0x100",
    "    行3 起始地址：0x6000   → 组索引 0x180",
    "    行4 起始地址：0x8000   → 组索引 0x000  (和行0冲突！)",
    "",
    "结果：多行竞争同一组的8个位置 → 频繁缓存驱逐"
])

# 第9页：问题e - 实验数据
add_table_slide(prs, "问题 e：不同矩阵规模的性能", [
    ["矩阵规模", "不转置(ms)", "转置后(ms)", "加速比", "备注"],
    ["1000", "114.69", "268.70", "0.43x", "普通"],
    ["1023", "218.71", "296.61", "0.74x", "2^-1 ✓"],
    ["1024", "189.27", "791.95", "0.24x", "2^幂 🔴"],
    ["1025", "235.27", "300.05", "0.78x", "2^+1 ✓"],
    ["2048", "1615.72", "7562.30", "0.21x", "2^幂 🔴"],
    ["2049", "1903.91", "2601.99", "0.73x", "2^+1 ✓"]
])

# 第10页：总结与建议
add_content_slide(prs, "总结与优化建议", [
    "问题 d 启示：",
    "    • 不是所有转置都能提升性能，需要分析访问模式",
    "    • 选择合适的循环顺序，事半功倍",
    "",
    "问题 e 启示：",
    "    • 避免使用 2 的幂次作为矩阵维度",
    "    • 使用填充（padding）改变数据布局，破坏冲突",
    "",
    "最优实践：",
    "    • 综合运用：kij顺序 + 合理的矩阵规模",
    "    • 编译优化：-O2/-O3 配合，性能提升最明显",
    "    • 性能测试：不同平台/编译器的优化策略可能不同"
])

# 保存PPT
output_path = "/Users/geyinhao/Desktop/studying/CS/2_1_de/矩阵乘法性能优化_问题d和e.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成：{output_path}")
