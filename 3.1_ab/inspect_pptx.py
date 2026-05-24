from pptx import Presentation
from pptx.util import Pt, Emu

prs = Presentation("report_2_1_a_v2.pptx")
slides = list(prs.slides)

for i, slide in enumerate(slides[:4]):
    print(f"\n=== Slide {i+1} ===")
    try:
        bg = slide.background.fill
        print(f"  BG: #{bg.fore_color.rgb}")
    except:
        print("  BG: none/theme")

    for shape in slide.shapes:
        x = Emu(shape.left).pt if shape.left else 0
        y = Emu(shape.top).pt if shape.top else 0
        w = Emu(shape.width).pt if shape.width else 0
        h = Emu(shape.height).pt if shape.height else 0
        print(f"  [{shape.shape_type}] {shape.name}: x={x:.0f} y={y:.0f} w={w:.0f} h={h:.0f}")
        try:
            fc = shape.fill.fore_color.rgb
            print(f"    fill=#{fc}")
        except:
            pass
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs[:2]:
                txt = para.text[:60]
                if not txt.strip():
                    continue
                for run in para.runs[:1]:
                    fn = run.font.name or "inherited"
                    fs = run.font.size
                    fs_pt = f"{fs.pt:.0f}pt" if fs else "inh"
                    bold = run.font.bold
                    try:
                        col = str(run.font.color.rgb)
                    except:
                        col = "theme"
                    print(f"    text={repr(txt[:40])} font={fn} sz={fs_pt} bold={bold} col=#{col}")
                    break
